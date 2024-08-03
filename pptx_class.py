from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR # center textbox
from pptx.enum.shapes import MSO_SHAPE #shapes
from pptx.dml.color import RGBColor #customize font
import locale
from datetime import datetime

class pptx():
    def __init__(self):
        self.prs = Presentation()

    def add_title_slide(self, title_text):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = str(title_text)

        locale.setlocale(locale.LC_TIME, 'es_ES.UTF-8')
        today = datetime.now()
        date_string = today.strftime('%A, %d de %B de %Y')

        subtitle.text = str(date_string)

    def add_song_slide(self, songlyrics):
        split_lyrics = songlyrics.splitlines()
        j = "\n"
        while len(split_lyrics) > 14:
            songpart = split_lyrics[0:13]
            self.part_add_song_slide(j.join(songpart))
            del split_lyrics[0:13]
        self.part_add_song_slide(j.join(split_lyrics))

    def part_add_song_slide(self, songlyrics):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide_width = self.prs.slide_width
        slide_height = self.prs.slide_height
        text_box_width = Inches(10)
        text_box_height = Inches(6)
        left = (slide_width - text_box_width) / 2
        top = (slide_height - text_box_height) / 2
        txBox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.text = str(songlyrics)

        # tf.paragraphs = None
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(32)
            paragraph.alignment = PP_ALIGN.CENTER

    def save_presentation(self, filename):
        self.prs.save(filename)


if 0 == 0:
    prs = pptx()
    prs.add_title_slide("Buttercup")
    prs.add_song_slide("Build me up buttercup \n Baby")
    prs.add_song_slide("13 going on 30")
    lyrics = """Purple rain, purple rain, purple drank, uh (yeah, ayy)
It's a purple party, all this purple on me, uh (purple)
She gon' pop it for me (pop), pop it for me (pop)
Purple rain, purple rain, purple drank, uh (yeah)
It's a purple party, all this purple on me, uh (purple)
She gon' pop it for me (pop), pop it for me (pop)

I love you the same
And I'm sippin' purple rain
And I hope you stay the same
And I hope you stay
I love you the same
And I'm sippin' purple rain
And I hope you stay the same
And I hope you stay

Four corners on the wall, I'm goin' crazy like I'm Wemby (Yeah)
Way deep down in her water, way too deep, finna break her levee (Sweet)
I'm standin' on the wings, lemon pepper (Sweet)
Ten-seater jet, this shit forever (Sweet)
Oh, what you hot for? Give me an encore
I got the Swisher on lil' shorty on the top floor
I'm steady runnin' up a check and, no, I'm not sore
It's Hardstone to the blood, bitch, it's hardcore

[Chorus: Future & Don Toliver]
Sippin' purple, sippin' purple rain
I'm sippin' purple, sippin' purple rain
Sippin' purple, yeah, I'm sippin' purple rain"""
    prs.add_song_slide(str(lyrics))
    prs.save_presentation("test buttercup 5.pptx")
