import tkinter as tk
from tkinter import ttk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class App():
    def __init__(self):
        self.root = tk.Tk()
        self.root.geometry('350x350')
        self.root.title('pptx Maker')
        self.mainframe = tk.Frame(self.root, background='pink')
        self.mainframe.pack(fill='both', expand=True)

        self.text = ttk.Label(self.mainframe, text='Praise the Lord', background='pink', font=("Brass Mono", 30))
        self.text.grid(row=0, column=0)

        self.set_text_field = ttk.Entry(self.mainframe)
        self.set_text_field.grid(row=1, column=0, pady=10, sticky='NWES')
        set_text_button = ttk.Button(self.mainframe, text='Set Text', command=self.lyric2pptx)
        set_text_button.grid(row=1, column=1, pady=10)

        options = ['red', 'blue']
        self.set_color_field = ttk.Combobox(self.mainframe, values=options)
        self.set_color_field.grid(row=2, column=0, sticky='NWES', pady=10)
        set_color_button = ttk.Button(self.mainframe, text='Set Color', command=self.lyric2pptx)
        set_color_button.grid(row=2, column=1, pady=10)

        self.root.mainloop()
        return

    def set_text(self):
        newtext = self.set_text_field.get()
        # self.text.config(text=newtext)

    def lyric2pptx(self):
        # Create a presentation object
        prs = Presentation()

        # Add a slide with a title and subtitle layout
        slide_layout = prs.slide_layouts[1]  # 1 is the layout index for title and content
        slide = prs.slides.add_slide(slide_layout)

        # Add a title and subtitle
        title = slide.shapes.title
        content = slide.placeholders[1]

        newtext = self.set_text_field.get()
        newcolor = self.set_color_field.get()
        title.text = newtext
        content.text = newcolor

        # Save the presentation
        prs.save("test7.pptx")


# if __name__ == '__lyric2pptx__':
App()
