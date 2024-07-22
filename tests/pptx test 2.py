from pptx import Presentation
from pptx.util import Inches, Pt #images and textbox
from pptx.enum.shapes import MSO_SHAPE #shapes
from pptx.dml.color import RGBColor #customize font
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = str(title_text)
    subtitle.text = str(subtitle_text)

# Create a presentation object
prs = Presentation()

# SLIDE 1: TITLE SLIDE
title_slide("Rainbow", "Buttercup")

# Add a slide with title and content
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Content Slide"
content.text = "This is a content slide with text."

slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank slide

# Define slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Define text box dimensions (larger than 1 inch)
text_box_width = Inches(6)
text_box_height = Inches(6)

# Calculate the position to center the text box on the slide
left = (slide_width - text_box_width) / 2
top = (slide_height - text_box_height) / 2

txBox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold \n Lorem ipsum lorem ipsum \n lorem ipsum"
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(40)
p.alignment = PP_ALIGN.CENTER


# Save the presentation
prs.save('example3.pptx')
