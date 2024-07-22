from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def create_presentation(text, filename):
    # Create a presentation object
    prs = Presentation()

    # Add a slide with a title and subtitle layout
    slide_layout = prs.slide_layouts[1]  # 1 is the layout index for title and content
    slide = prs.slides.add_slide(slide_layout)

    # Add a title and subtitle
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Title Example"
    content.text = text

    # Save the presentation
    prs.save(filename)

#-----------------------------------------------------------------------------------

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Define slide dimensions
slide_width = prs.slide_width
slide_height = prs.slide_height

# Define text box dimensions (larger than 1 inch)
text_box_width = Inches(6)
text_box_height = Inches(6)

# Calculate the position to center the text box on the slide
left = (slide_width - text_box_width) / 2
top = (slide_height - text_box_height) / 2

txBox = slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.add_paragraph()
p.text = "This is a second paragraph that's bold \n Lorem ipsum lorem ipsum \n lorem ipsum"
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(40)
p.alignment = PP_ALIGN.CENTER


prs.save('test5.pptx')

